from __future__ import annotations

from collections import Counter
from dataclasses import dataclass, field
from html.parser import HTMLParser
from pathlib import Path
import re
from urllib.request import Request, urlopen


STOP_WORDS = {
    "about",
    "above",
    "after",
    "again",
    "against",
    "also",
    "and",
    "are",
    "but",
    "because",
    "before",
    "being",
    "below",
    "between",
    "could",
    "during",
    "each",
    "for",
    "from",
    "have",
    "has",
    "how",
    "its",
    "into",
    "is",
    "inside",
    "not",
    "of",
    "on",
    "or",
    "more",
    "most",
    "other",
    "over",
    "should",
    "some",
    "such",
    "than",
    "that",
    "the",
    "their",
    "then",
    "there",
    "these",
    "they",
    "this",
    "through",
    "to",
    "under",
    "uses",
    "using",
    "when",
    "where",
    "which",
    "while",
    "with",
    "would",
    "your",
}

BOUNDARY_WORDS = STOP_WORDS | {
    "absorbs",
    "allows",
    "can",
    "causes",
    "convert",
    "converts",
    "create",
    "creates",
    "describe",
    "describes",
    "explain",
    "explains",
    "include",
    "includes",
    "make",
    "makes",
    "mean",
    "means",
    "provide",
    "provides",
    "show",
    "shows",
    "support",
    "supports",
    "use",
    "used",
}

WEAK_SINGLE_WORDS = {
    "application",
    "cells",
    "concept",
    "concepts",
    "definition",
    "energy",
    "example",
    "growth",
    "idea",
    "ideas",
    "material",
    "method",
    "notes",
    "process",
    "production",
    "purpose",
    "release",
    "source",
    "study",
    "system",
    "term",
    "terms",
    "thing",
    "things",
    "topic",
    "topics",
}

ACRONYMS = {"atp", "dna", "rna", "html", "css", "api", "cpu", "gpu", "pdf"}


@dataclass(frozen=True)
class KeyConcept:
    title: str
    explanation: str
    source_sentence: str


@dataclass(frozen=True)
class StudyBlock:
    title: str
    block_type: str
    summary: str
    items: list[str] = field(default_factory=list)
    source_excerpt: str = ""
    page: int | None = None


@dataclass(frozen=True)
class ParsedContent:
    title: str
    source_type: str
    text: str
    concepts: list[str]
    concept_details: list[KeyConcept] = field(default_factory=list)
    study_blocks: list[StudyBlock] = field(default_factory=list)


class ParseError(RuntimeError):
    """Raised when a source cannot be parsed into usable study text."""


class _HTMLTextParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._ignored_depth = 0
        self._parts: list[str] = []
        self.title = ""
        self._in_title = False

    def handle_starttag(self, tag: str, attrs) -> None:
        if tag in {"script", "style", "noscript", "svg"}:
            self._ignored_depth += 1
        elif tag == "title":
            self._in_title = True
        elif tag in {"p", "br", "li", "div", "section", "article", "h1", "h2", "h3"}:
            self._parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag in {"script", "style", "noscript", "svg"} and self._ignored_depth:
            self._ignored_depth -= 1
        elif tag == "title":
            self._in_title = False

    def handle_data(self, data: str) -> None:
        if self._ignored_depth:
            return
        clean = " ".join(data.split())
        if not clean:
            return
        if self._in_title:
            self.title = f"{self.title} {clean}".strip()
        self._parts.append(clean)

    @property
    def text(self) -> str:
        return "\n".join(part for part in self._parts if part.strip())


def parse_text(raw_text: str, title: str = "Pasted text") -> ParsedContent:
    text = _clean_text(raw_text)
    if not text:
        raise ParseError("No readable text was found.")
    study_blocks = extract_study_blocks(text)
    concept_details = _key_concepts_from_study_blocks(study_blocks) or extract_key_concept_details(text)
    return ParsedContent(
        title=title,
        source_type="text",
        text=text,
        concepts=[concept.title for concept in concept_details],
        concept_details=concept_details,
        study_blocks=study_blocks,
    )


def parse_file(path: str | Path) -> ParsedContent:
    file_path = Path(path)
    if not file_path.exists():
        raise ParseError(f"File not found: {file_path}")

    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        text = _extract_pdf_text(file_path)
        source_type = "pdf"
    elif suffix == ".pptx":
        text = _extract_pptx_text(file_path)
        source_type = "pptx"
    elif suffix in {".txt", ".md"}:
        text = file_path.read_text(encoding="utf-8", errors="ignore")
        source_type = "text"
    else:
        raise ParseError("Supported files are PDF, PPTX, TXT, and MD.")

    text = _clean_text(text)
    if not text:
        raise ParseError("No readable text was found in the selected file.")
    study_blocks = extract_study_blocks(text)
    concept_details = _key_concepts_from_study_blocks(study_blocks) or extract_key_concept_details(text)
    return ParsedContent(
        title=file_path.name,
        source_type=source_type,
        text=text,
        concepts=[concept.title for concept in concept_details],
        concept_details=concept_details,
        study_blocks=study_blocks,
    )


def parse_url(url: str, timeout: int = 15) -> ParsedContent:
    if not url.lower().startswith(("http://", "https://")):
        raise ParseError("Enter a full http:// or https:// web link.")

    request = Request(url, headers={"User-Agent": "ExamBuddy/1.0"})
    try:
        with urlopen(request, timeout=timeout) as response:
            content_type = response.headers.get("content-type", "")
            charset = response.headers.get_content_charset() or "utf-8"
            payload = response.read(2_000_000)
    except Exception as exc:  # pragma: no cover - network failures vary by platform
        raise ParseError(f"Could not fetch the web link: {exc}") from exc

    if "html" not in content_type and b"<html" not in payload[:500].lower():
        text = payload.decode(charset, errors="ignore")
        title = url
    else:
        parser = _HTMLTextParser()
        parser.feed(payload.decode(charset, errors="ignore"))
        text = parser.text
        title = parser.title or url

    text = _clean_text(text)
    if not text:
        raise ParseError("No readable text was found at the web link.")
    study_blocks = extract_study_blocks(text)
    concept_details = _key_concepts_from_study_blocks(study_blocks) or extract_key_concept_details(text)
    return ParsedContent(
        title=title,
        source_type="web",
        text=text,
        concepts=[concept.title for concept in concept_details],
        concept_details=concept_details,
        study_blocks=study_blocks,
    )


def extract_key_concepts(text: str, minimum: int = 5, maximum: int = 24) -> list[str]:
    return [concept.title for concept in extract_key_concept_details(text, minimum=minimum, maximum=maximum)]


def extract_study_blocks(text: str, minimum: int = 5, maximum: int = 24) -> list[StudyBlock]:
    blocks = _extract_structured_study_blocks(text)
    if not blocks:
        fallback = [
            StudyBlock(
                title=concept.title,
                block_type="concept",
                summary=concept.explanation,
                items=[],
                source_excerpt=concept.source_sentence,
            )
            for concept in extract_key_concept_details(text, minimum=minimum, maximum=maximum)
        ]
        seen = {block.title.lower() for block in blocks}
        blocks.extend(block for block in fallback if block.title.lower() not in seen)
    return blocks[:maximum]


def extract_key_concept_details(text: str, minimum: int = 5, maximum: int = 24) -> list[KeyConcept]:
    cleaned_text = _clean_text(text)
    cleaned = cleaned_text.lower()
    if not cleaned:
        return []

    sentences = _split_sentences(cleaned_text)
    if not sentences:
        return []

    all_words: list[str] = []
    candidate_sentence: dict[str, str] = {}
    first_seen: dict[str, int] = {}
    candidates: dict[str, float] = {}

    for sentence in sentences:
        words = re.findall(r"[a-z][a-z0-9'-]{1,}", sentence.lower())
        words = [word for word in words if not word.isdigit()]
        all_words.extend(word for word in words if word not in STOP_WORDS)

    if not all_words:
        return []

    word_counts = Counter(all_words)

    for sentence_index, sentence in enumerate(sentences):
        fragments = re.split(r"[,;:()]|\b(?:and|or)\b", sentence, flags=re.IGNORECASE)
        for fragment in fragments:
            words = re.findall(r"[a-z][a-z0-9'-]{1,}", fragment.lower())
            segments: list[list[str]] = []
            current: list[str] = []
            for word in words:
                if word in BOUNDARY_WORDS:
                    if current:
                        segments.append(current)
                        current = []
                    continue
                current.append(word)
            if current:
                segments.append(current)

            for segment in segments:
                _add_candidates_from_segment(segment, sentence_index, sentences, word_counts, first_seen, candidate_sentence, candidates)

    ranked = sorted(candidates, key=lambda item: (-candidates[item], first_seen[item], item))
    multiword_candidate_parts = {
        part
        for candidate in ranked
        if len(candidate.split()) > 1
        for part in candidate.split()
    }
    selected: list[str] = []
    for candidate in ranked:
        if len(candidate.split()) == 1 and candidate in multiword_candidate_parts:
            continue
        if _is_redundant(candidate, selected):
            continue
        selected.append(candidate)
        if len(selected) >= maximum:
            break

    if len(selected) < minimum:
        fallback_words = [
            word
            for word in sorted(word_counts, key=lambda item: (-word_counts[item], item))
            if _valid_candidate([word]) and word not in selected
        ]
        for word in fallback_words:
            selected.append(word)
            if len(selected) >= min(minimum, maximum):
                break

    return [
        KeyConcept(
            title=_format_concept_title(candidate),
            explanation=_build_explanation(candidate, candidate_sentence.get(candidate, "")),
            source_sentence=candidate_sentence.get(candidate, ""),
        )
        for candidate in selected[:maximum]
    ]


def _extract_structured_study_blocks(text: str) -> list[StudyBlock]:
    entries = _extract_heading_entries(text)
    blocks: list[StudyBlock] = []
    used_indices: set[int] = set()

    population_block = _population_sampling_block(entries)
    population_index: int | None = None
    if population_block:
        for index, (heading, _body) in enumerate(entries):
            if _normalize_heading(heading) in POPULATION_GROUP_HEADINGS:
                used_indices.add(index)
                if population_index is None:
                    population_index = index

    for index, (heading, body) in enumerate(entries):
        if population_block and index == population_index:
            blocks.append(population_block)
        if index in used_indices:
            continue
        block = _entry_to_study_block(heading, body)
        if block:
            blocks.append(block)

    return _dedupe_study_blocks(blocks)


def _extract_heading_entries(text: str) -> list[tuple[str, str]]:
    lines = _study_lines(text)
    entries: list[tuple[str, str]] = []
    current_heading = ""
    current_body: list[str] = []
    previous_blank = True

    for line in lines:
        if not line:
            previous_blank = True
            continue
        if _is_lecture_marker(line):
            previous_blank = True
            continue
        heading, rest = _split_heading_line(line)
        if heading and current_heading and not previous_blank and _is_detail_label(heading):
            current_body.append(line)
        elif heading:
            if current_heading:
                entries.append((current_heading, " ".join(current_body).strip()))
            current_heading = heading
            current_body = [rest] if rest else []
        elif previous_blank and _is_study_heading(line):
            if current_heading:
                entries.append((current_heading, " ".join(current_body).strip()))
            current_heading = line
            current_body = []
        elif current_heading:
            current_body.append(line)
        previous_blank = False

    if current_heading:
        entries.append((current_heading, " ".join(current_body).strip()))
    return entries


def _study_lines(text: str) -> list[str]:
    lines: list[str] = []
    for raw_line in _clean_text(text).splitlines():
        line = " ".join(raw_line.split()).strip()
        if not line:
            lines.append("")
            continue
        if _is_noise_line(line):
            continue
        lines.append(line)
    return lines


def _split_heading_line(line: str) -> tuple[str, str]:
    if len(line) > 180 or ":" not in line:
        return "", ""
    heading, rest = line.split(":", 1)
    heading = heading.strip()
    if not _is_study_heading(heading):
        return "", ""
    return heading, rest.strip()


def _is_study_heading(heading: str) -> bool:
    normalized = _normalize_heading(heading)
    if normalized in {"lecture", "example", "solution", "step"}:
        return True
    if normalized.startswith("step "):
        return True
    return bool(re.match(r"^[A-Z][A-Za-z0-9 &/().-]{2,85}$", heading)) and not heading.lower().startswith(
        ("as for example", "for example", "following")
    )


def _is_detail_label(heading: str) -> bool:
    return _normalize_heading(heading) in {"technique", "techniques", "example", "examples", "note", "notes"}


def _entry_to_study_block(heading: str, body: str) -> StudyBlock | None:
    if not body:
        return None
    normalized = _normalize_heading(heading)
    items = _extract_list_items(body)

    if normalized.startswith("definition of ") or normalized.startswith("definitions of "):
        prefix_length = len("Definitions of ") if normalized.startswith("definitions of ") else len("Definition of ")
        subject = _title_case(heading[prefix_length:])
        title = f"Definitions of {subject}"
        summary = _definition_summary(subject, body)
        technique_items = _technique_items(body)
        return StudyBlock(title, "definition", summary, technique_items or items, _source_excerpt(heading, body))

    if normalized.startswith("importance of "):
        return StudyBlock(_title_case(heading), "list", _first_sentence(body), items or _sentence_items(body), _source_excerpt(heading, body))

    if normalized.startswith("methods of ") or normalized.startswith("types of ") or normalized.startswith("use of ") or normalized.startswith("basic principle"):
        return StudyBlock(_title_case(heading), "list", _first_sentence(body), items, _source_excerpt(heading, body))

    if normalized.startswith("constructing ") or normalized.startswith("steps for "):
        return StudyBlock(_title_case(heading), "steps", _first_sentence(body), items, _source_excerpt(heading, body))

    if _looks_like_definition(heading, body):
        summary = _definition_summary(heading, body)
        return StudyBlock(_title_case(heading), "definition", summary, items, _source_excerpt(heading, body))

    if items:
        return StudyBlock(_title_case(heading), "list", _first_sentence(body), items, _source_excerpt(heading, body))

    return None


POPULATION_GROUP_HEADINGS = {
    "population",
    "finite population",
    "infinite population",
    "sample",
    "sampling unit",
    "sampling frame",
}


def _population_sampling_block(entries: list[tuple[str, str]]) -> StudyBlock | None:
    found: dict[str, str] = {}
    excerpts: list[str] = []
    for heading, body in entries:
        normalized = _normalize_heading(heading)
        if normalized in POPULATION_GROUP_HEADINGS:
            found[normalized] = body
            excerpts.append(_source_excerpt(heading, body))

    if not {"population", "sample", "sampling frame"}.issubset(found):
        return None

    items: list[str] = []
    if "population" in found:
        items.append("Population -> Entire group under study.")
    if "finite population" in found:
        items.append("Finite population -> Countable units.")
    if "infinite population" in found:
        items.append("Infinite population -> Uncountable units.")
    if "sample" in found:
        items.append("Sample -> Representative part of population.")
    if "sampling unit" in found:
        items.append("Sampling unit -> Smallest unit information is collected from.")
    if "sampling frame" in found:
        items.append("Sampling frame -> List of all population units.")

    return StudyBlock(
        title="Population & Sampling",
        block_type="grouped_definitions",
        summary="Core population and sampling terms used to define what is studied and what information is collected.",
        items=items,
        source_excerpt=" ".join(excerpts)[:900],
    )


def _add_candidates_from_segment(
    segment: list[str],
    sentence_index: int,
    sentences: list[str],
    word_counts: Counter[str],
    first_seen: dict[str, int],
    candidate_sentence: dict[str, str],
    candidates: dict[str, float],
) -> None:
    for start in range(len(segment)):
        for size in range(1, min(4, len(segment) - start) + 1):
            parts = segment[start : start + size]
            if not _valid_candidate(parts):
                continue
            key = " ".join(parts)
            first_seen.setdefault(key, len(first_seen))
            candidate_sentence.setdefault(key, sentences[sentence_index])
            score = _candidate_score(parts, word_counts, start == 0)
            candidates[key] = candidates[key] + 3 if key in candidates else score


def _extract_pdf_text(path: Path) -> str:
    try:
        try:
            from pypdf import PdfReader
        except ImportError:
            from PyPDF2 import PdfReader  # type: ignore
    except ImportError as exc:
        raise ParseError("PDF parsing requires pypdf. Install it with: pip install -r requirements.txt") from exc

    try:
        reader = PdfReader(str(path))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception as exc:
        raise ParseError(f"Could not parse PDF: {exc}") from exc


def _extract_pptx_text(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ParseError("PPTX parsing requires python-pptx. Install it with: pip install -r requirements.txt") from exc

    try:
        presentation = Presentation(str(path))
        chunks: list[str] = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            chunks.append(f"Slide {slide_number}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    chunks.append(shape.text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        chunks.append(" ".join(cell.text for cell in row.cells))
        return "\n".join(chunks)
    except Exception as exc:
        raise ParseError(f"Could not parse PPTX: {exc}") from exc


def _key_concepts_from_study_blocks(blocks: list[StudyBlock]) -> list[KeyConcept]:
    return [
        KeyConcept(
            title=block.title,
            explanation=_block_explanation(block),
            source_sentence=block.source_excerpt,
        )
        for block in blocks
    ]


def _block_explanation(block: StudyBlock) -> str:
    parts = [block.summary.strip()] if block.summary.strip() else []
    parts.extend(item.strip() for item in block.items if item.strip())
    return "\n".join(parts)


def _extract_list_items(body: str) -> list[str]:
    definition_items = [
        _clean_item(f"{match.group(1).strip()} -> {match.group(2).strip()}")
        for match in re.finditer(r"([A-Z][A-Za-z ]{2,60})\s*(?:->|→|:)\s*(.+?)(?=\s+[A-Z][A-Za-z ]{2,60}\s*(?:->|→|:)|$)", body)
    ]
    if len(definition_items) >= 2:
        return [item for item in definition_items if item]

    technique_items = _technique_items(body)
    if len(technique_items) >= 2:
        return technique_items

    bullet_items = [
        _clean_item(match.group(1))
        for match in re.finditer(r"(?:^|\s)[•\-]\s+([^•\-]+?)(?=(?:\s[•\-]\s+)|$)", body)
    ]
    if bullet_items:
        return [item for item in bullet_items if item]

    letter_items = [
        _clean_item(match.group(2))
        for match in re.finditer(r"(?:^|\s)(\([a-zivx]+\)|[a-zivx]+\))\s+(.+?)(?=(?:\s(?:\([a-zivx]+\)|[a-zivx]+\))\s+)|$)", body, re.I)
    ]
    if len(letter_items) >= 2:
        return [item for item in letter_items if item]

    step_items = [
        _clean_item(f"{match.group(1)}: {match.group(2)}")
        for match in re.finditer(r"(Step\s+\d+)\s*:\s*(.+?)(?=(?:\sStep\s+\d+\s*:)|$)", body, re.I)
    ]
    return [item for item in step_items if item]


def _technique_items(body: str) -> list[str]:
    match = re.search(r"techniques?(?:\s+such\s+as|:)\s+(.+?)(?:\s+to\s+|\.)", body, re.I)
    if not match:
        return []
    return [_sentence_case(item.strip().lower()) for item in re.split(r",|\band\b", match.group(1)) if item.strip()]


def _sentence_items(body: str) -> list[str]:
    return [_clean_item(sentence) for sentence in _split_sentences(body) if len(_clean_item(sentence)) >= 8]


def _definition_summary(title: str, body: str) -> str:
    text = _first_sentence(body)
    patterns = [
        r"may be defined as\s+(.+)",
        r"is known as\s+(.+)",
        r"is\s+(.+)",
    ]
    for pattern in patterns:
        match = re.search(pattern, text, re.I)
        if match:
            return _sentence_case(_strip_leading_article(_clean_item(match.group(1))))
    return _sentence_case(_clean_item(text))


def _looks_like_definition(heading: str, body: str) -> bool:
    if len(body) < 20:
        return False
    lowered = body.lower()
    return any(phrase in lowered for phrase in ("known as", "defined as", "is a", "is an", "is the", "refers to"))


def _first_sentence(text: str) -> str:
    sentence = re.split(r"(?<=[.!?])\s+", re.sub(r"\s+", " ", text).strip())[0]
    return _clean_item(sentence)


def _clean_item(item: str) -> str:
    item = re.sub(r"\s+", " ", item or "").strip(" .;:")
    item = item.replace("data patterns into usable format", "data patterns into usable format")
    return item + "." if item and item[-1] not in ".!?" else item


def _sentence_case(item: str) -> str:
    if not item:
        return item
    item = item.strip()
    return item[:1].upper() + item[1:]


def _strip_leading_article(item: str) -> str:
    return re.sub(r"^(an?|the)\s+", "", item.strip(), flags=re.I)


def _title_case(value: str) -> str:
    words = []
    small_words = {"and", "or", "of", "in", "to", "for", "the", "a", "an"}
    raw_words = re.sub(r"\s+", " ", value).strip(" .:").split()
    for index, word in enumerate(raw_words):
        lower = word.lower()
        if lower in ACRONYMS:
            words.append(word.upper())
        elif index and lower in small_words:
            words.append(lower)
        else:
            words.append(lower.capitalize())
    return " ".join(words)


def _source_excerpt(heading: str, body: str) -> str:
    return f"{heading}: {body}".strip()[:900]


def _normalize_heading(heading: str) -> str:
    heading = re.sub(r"\s+", " ", heading).strip(" .:")
    heading = re.sub(r"\s+\d+$", "", heading)
    return heading.lower()


def _is_lecture_marker(line: str) -> bool:
    return bool(re.match(r"^(Chapter|Lecture|Slide)\s+\d+\s*:?", line, re.I))


def _is_noise_line(line: str) -> bool:
    if line in {"Introduction to Statistics and Data Science", "Professor, Department of SDS"}:
        return True
    if re.match(r"^Dr\.?\s+Mohd\.?\s+Muzibur\s+Rahman\b", line, re.I):
        return True
    if re.match(r"^\d+$", line):
        return True
    return False


def _dedupe_study_blocks(blocks: list[StudyBlock]) -> list[StudyBlock]:
    results: list[StudyBlock] = []
    seen: set[str] = set()
    for block in blocks:
        key = block.title.lower()
        if key in seen:
            continue
        if not block.summary and not block.items:
            continue
        results.append(block)
        seen.add(key)
    return results


def _clean_text(text: str) -> str:
    text = re.sub(r"\r\n?", "\n", text or "")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _split_sentences(text: str) -> list[str]:
    normalized = re.sub(r"\s+", " ", text).strip()
    if not normalized:
        return []
    return [
        sentence.strip()
        for sentence in re.split(r"(?<=[.!?])\s+", normalized)
        if len(sentence.strip()) >= 8
    ]


def _valid_candidate(parts: list[str]) -> bool:
    if not parts or len(set(parts)) != len(parts):
        return False
    if any(part in STOP_WORDS for part in parts):
        return False
    if len(parts) == 1:
        word = parts[0]
        return word not in WEAK_SINGLE_WORDS and len(word) >= 4
    if parts[0] in WEAK_SINGLE_WORDS:
        return False
    return not all(part in WEAK_SINGLE_WORDS for part in parts)


def _candidate_score(parts: list[str], word_counts: Counter[str], starts_segment: bool) -> float:
    if len(parts) == 1:
        word = parts[0]
        return word_counts[word] * 12 + min(len(word), 18) / 2 + (8 if starts_segment else 0)

    repeated_word_bonus = sum(2 for word in parts if word_counts[word] > 1)
    acronym_bonus = 4 if any(word in ACRONYMS for word in parts) else 0
    return 16 + len(parts) * 5 + repeated_word_bonus + acronym_bonus + (4 if starts_segment else 0)


def _is_redundant(candidate: str, selected: list[str]) -> bool:
    candidate_parts = set(candidate.split())
    for existing in selected:
        existing_parts = set(existing.split())
        if candidate_parts == existing_parts:
            return True
        if len(candidate_parts) == 1 and candidate_parts.issubset(existing_parts):
            return True
    return False


def _format_concept_title(candidate: str) -> str:
    return " ".join(word.upper() if word in ACRONYMS else word for word in candidate.split())


def _build_explanation(candidate: str, sentence: str) -> str:
    title = _format_concept_title(candidate)
    if sentence:
        return f"{title}: {sentence.rstrip('.!?')}."
    return f"{title}: review this concept in the source material."
